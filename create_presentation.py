import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import io

# Load the data
df = pd.read_csv('module1_project_dashboard/data/sample_data.csv')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 118, 210)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_chart_slide(prs, title, fig):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(25, 118, 210)
    
    img_stream = io.BytesIO()
    fig.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.3), width=Inches(9))
    plt.close(fig)

# Slide 1: Title Slide
add_title_slide(prs, 'Market Trend Analysis', 'Strategic Industry Growth Assessment')

# Slide 2: Business Objective
add_content_slide(prs, 'Business Objective', [
    ' Conduct a comprehensive market trend analysis to identify high-growth industry segments',
    ' Analyze multi-year revenue, market size, and growth rate data across key industries',
    ' Provide data-driven insights to support strategic product portfolio decisions'
])

# Slide 3: Target User & Strategic Use
add_content_slide(prs, 'Target User & Strategic Application', [
    'Target User: Product Leadership & Strategic Planning Teams',
    ' Identify high-growth segments for strategic prioritization',
    ' Make informed decisions on which products receive funding and which are sunsetted',
    ' Allocate resources based on industry growth trajectories and market opportunities'
])

# Slide 4: Business Value Proposition
add_content_slide(prs, 'Business Value Proposition', [
    'Maximize Resource Efficiency: Direct investment towards highest-growth segments',
    ' Minimize opportunity cost by reducing investment in low-growth or declining segments',
    ' Optimize capital allocation to achieve maximum ROI',
    ' Data-driven approach reduces strategic risk and improves decision outcomes'
])

# Slide 5: Growth Rate by Industry (Latest Year)
latest_year = df[df['year'] == df['year'].max()]
industry_growth = latest_year.groupby('industry')['growth_rate'].mean().sort_values(ascending=False)

fig, ax = plt.subplots(figsize=(10, 6))
colors = ['#1976D2' if x >= 0.25 else '#FFA726' if x >= 0.15 else '#66BB6A' for x in industry_growth.values]
bars = ax.barh(industry_growth.index, industry_growth.values, color=colors)
ax.set_xlabel('Average Growth Rate (2022)', fontsize=12, fontweight='bold')
ax.set_title('Industry Growth Rate Comparison', fontsize=14, fontweight='bold', pad=20)
ax.set_xlim(0, max(industry_growth.values) * 1.1)

for i, bar in enumerate(bars):
    width = bar.get_width()
    ax.text(width, bar.get_y() + bar.get_height()/2, f'{width:.1%}', 
            ha='left', va='center', fontweight='bold', fontsize=11)

plt.tight_layout()
add_chart_slide(prs, 'Industry Growth Rates (2022)', fig)

# Slide 6: Revenue Trends by Industry
fig, ax = plt.subplots(figsize=(10, 6))
for industry in df['industry'].unique():
    industry_data = df[df['industry'] == industry].sort_values('year')
    ax.plot(industry_data['year'], industry_data['revenue'], marker='o', linewidth=2.5, 
            label=industry, markersize=6)

ax.set_xlabel('Year', fontsize=12, fontweight='bold')
ax.set_ylabel('Revenue (Millions)', fontsize=12, fontweight='bold')
ax.set_title('Revenue Trend Analysis by Industry (2018-2022)', fontsize=14, fontweight='bold', pad=20)
ax.legend(loc='upper left', fontsize=10)
ax.grid(True, alpha=0.3)
plt.tight_layout()
add_chart_slide(prs, 'Revenue Trends by Industry', fig)

# Slide 7: Market Size Growth Potential
latest_year_data = df[df['year'] == df['year'].max()].groupby('industry')['market_size'].sum()
earliest_year_data = df[df['year'] == df['year'].min()].groupby('industry')['market_size'].sum()
growth_potential = ((latest_year_data - earliest_year_data) / earliest_year_data * 100).sort_values(ascending=False)

fig, ax = plt.subplots(figsize=(10, 6))
colors = ['#1976D2' if x >= 150 else '#FFA726' if x >= 100 else '#66BB6A' for x in growth_potential.values]
bars = ax.barh(growth_potential.index, growth_potential.values, color=colors)
ax.set_xlabel('Market Size Growth 2018-2022 (%)', fontsize=12, fontweight='bold')
ax.set_title('Total Market Size Expansion Potential', fontsize=14, fontweight='bold', pad=20)
ax.set_xlim(0, max(growth_potential.values) * 1.1)

for i, bar in enumerate(bars):
    width = bar.get_width()
    ax.text(width, bar.get_y() + bar.get_height()/2, f'{width:.0f}%', 
            ha='left', va='center', fontweight='bold', fontsize=11)

plt.tight_layout()
add_chart_slide(prs, 'Market Size Expansion (2018-2022)', fig)

# Slide 8: Strategic Recommendations
add_content_slide(prs, 'Strategic Recommendations', [
    ' High-Growth Priority (Growth Rate >= 30%):',
    '   Focus investment on Enterprise AI and E-Commerce segments',
    ' Sunset Candidates (Growth Rate <= 5%):',
    '   Consider strategic review of Healthcare Devices segment',
    ' Balanced Portfolio: Maintain presence in stable segments while scaling high-growth opportunities'
])

# Slide 9: Summary & Next Steps
add_content_slide(prs, 'Summary & Next Steps', [
    ' Identified Enterprise AI and E-Commerce as fastest-growing segments (40%+ growth)',
    ' Quantified market expansion potential across all industries',
    ' Recommend reallocating 30-40% of resources to high-growth segments',
    ' Schedule board review for product portfolio optimization',
    ' Establish quarterly monitoring dashboard for market trend updates'
])

# Save presentation
prs.save('Market_Trend_Analysis_Presentation.pptx')
print(' PowerPoint presentation created successfully: Market_Trend_Analysis_Presentation.pptx')
